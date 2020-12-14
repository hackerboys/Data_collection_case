# encoding:utf-8
"""
作者：宋哈哈
代码功能：
    删除PPT最后一页有logo的。方便与上传其他网站等
使用说明：
    请自行更改保存的PPT路径和另存的路径 和 index_down_ppt.py 的路径相连。

"""



from pptx import Presentation
def pptxde(pptfilepath,pptsavepath):
    """
    删除PPT末尾倒数第一第二或倒数第一页
    :param pptpath: PPT文件路径
    :param pptname: ppt文件名称
    :return:
    """

    pptlist = []
    prs = Presentation(r'%s'%pptfilepath)
    for p in prs.slides:
        for s in p.shapes:
            if s.has_text_frame:
                text_frame = s.text_frame
                for t in text_frame.paragraphs:
                    pptlist.append(t.text)

    strlist = "".join(pptlist)
    if "baotu" in str(strlist):
        sli = list(prs.slides._sldIdLst)
        prs.slides._sldIdLst.remove(sli[-1])
        prs.slides._sldIdLst.remove(sli[-2])
        prs.save(r'%s'%pptsavepath)
    else:
        sli = list(prs.slides._sldIdLst)
        prs.slides._sldIdLst.remove(sli[-1])
        prs.save(r'%s'%pptsavepath)


if __name__ == '__main__':
    pptxde()